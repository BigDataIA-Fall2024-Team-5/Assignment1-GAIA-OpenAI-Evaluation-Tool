import os
import pandas as pd
import json
from PIL import Image
import zipfile
from docx import Document
from PyPDF2 import PdfReader

def preprocess_file(file_path):
    """Preprocess a file based on its extension and return relevant information."""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.txt':
        return preprocess_txt(file_path)
    elif file_extension == '.csv':
        return preprocess_csv(file_path)
    elif file_extension == '.xlsx':
        return preprocess_xlsx(file_path)
    elif file_extension in ['.jpg', '.png']:
        return preprocess_image(file_path)
    elif file_extension == '.jsonld':
        return preprocess_jsonld(file_path)
    elif file_extension == '.zip':
        return preprocess_zip(file_path)
    elif file_extension == '.docx':
        return preprocess_docx(file_path)
    elif file_extension == '.pdf':
        return preprocess_pdf(file_path)
    elif file_extension == '.py':
        return preprocess_py(file_path)
    elif file_extension == '.pptx':
        return preprocess_pptx(file_path)
    elif file_extension == '.pdb':
        return preprocess_pdb(file_path)
    elif file_extension == '.mp3':
        return preprocess_mp3(file_path)
    else:
        return f"Unsupported file type: {file_extension}"

def preprocess_txt(file_path):
    """Preprocess a .txt file by reading and returning its content."""
    return read_file_content(file_path)

def preprocess_csv(file_path):
    """Preprocess a .csv file by loading and returning its full content or as much as possible if too large."""
    try:
        df = pd.read_csv(file_path)
        content = df.to_string(index=False)
        
        # Adjust based on the token limit (e.g., 16000 characters)
        if len(content) > 16000:
            content = content[:16000]  # Truncate to fit the limit
            
        return {"content": content}
    except Exception as e:
        return f"Error processing CSV file: {e}"

def preprocess_xlsx(file_path):
    """Preprocess an .xlsx file by loading and returning its full content or as much as possible if too large."""
    try:
        df = pd.read_excel(file_path, sheet_name=0)
        content = df.to_string(index=False)
        
        # Adjust based on the token limit
        if len(content) > 16000:
            content = content[:16000]  # Truncate to fit the limit
            
        return {"content": content}
    except Exception as e:
        return f"Error processing XLSX file: {e}"

def preprocess_image(file_path):
    """Preprocess an image file by returning its size and mode."""
    try:
        with Image.open(file_path) as img:
            return {"info": {"size": img.size, "mode": img.mode}}
    except Exception as e:
        return f"Error processing image file: {e}"

def preprocess_jsonld(file_path):
    """Preprocess a .jsonld file by loading and returning its content."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        return {"content": json.dumps(data, indent=2)}
    except Exception as e:
        return f"Error processing JSON-LD file: {e}"

def preprocess_zip(file_path):
    """Preprocess a .zip file by returning a list of its contents."""
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            return {"content": zip_ref.namelist()}
    except Exception as e:
        return f"Error processing ZIP file: {e}"

def preprocess_docx(file_path):
    """Preprocess a .docx file by reading and returning its content."""
    try:
        doc = Document(file_path)
        full_text = []
        for paragraph in doc.paragraphs:
            full_text.append(paragraph.text)
        return '\n'.join(full_text)
    except Exception as e:
        return f"Error processing DOCX file: {e}"

def preprocess_pdf(file_path):
    """Preprocess a .pdf file by extracting and returning its text."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        
        # Adjust based on the token limit
        if len(text) > 16000:
            text = text[:16000]  # Truncate to fit the limit
            
        return text
    except Exception as e:
        return f"Error processing PDF file: {e}"

def preprocess_py(file_path):
    """Preprocess a .py file by returning its content."""
    return read_file_content(file_path)

def preprocess_pptx(file_path):
    """Preprocess a .pptx file by extracting and returning slide content."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)
        return '\n'.join(slides_text)
    except Exception as e:
        return f"Error processing PPTX file: {e}"

def preprocess_pdb(file_path):
    """Preprocess a .pdb file by returning its content."""
    return read_file_content(file_path)

def preprocess_mp3(file_path):
    """Preprocess an .mp3 file by returning its metadata."""
    try:
        from mutagen.mp3 import MP3
        audio = MP3(file_path)
        return {"duration": audio.info.length, "bitrate": audio.info.bitrate}
    except Exception as e:
        return f"Error processing MP3 file: {e}"

def read_file_content(file_path):
    """Read content from a file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error reading file {file_path}: {e}"
